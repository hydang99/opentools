# source code: https://github.com/inclusionAI/AWorld/blob/main/examples/gaia/mcp_collections/documents/mspptx.py
import json, os, sys, time, traceback, zipfile
from pathlib import Path
from typing import Any, Literal, Optional
from openpyxl import load_workbook
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__),'..', '..', '..')))
from opentools.core.base import BaseTool
from pydantic import Field, BaseModel
from pydantic.fields import FieldInfo
from pptx import Presentation
from pptx.presentation import Presentation as PresentationType

class DocumentMetadata(BaseModel):
    """Metadata for extracted PowerPoint document content."""
    file_name: str
    file_size: int
    file_type: str
    absolute_path: str
    page_count: Optional[int] = None
    processing_time: float
    extracted_images: list[str] = []
    extracted_media: list[dict[str, str]] = []
    output_format: str
    ocr_applied: bool
    extracted_text_file_path: Optional[str] = None


class Pptx_Extraction_Tool(BaseTool):
    # Default args for `opentools test Pptx_Extraction_Tool` (uses test_file/data.json)
    DEFAULT_TEST_ARGS = {
        "tool_test": "pptx_extraction",
        "file_location": "pptx_extraction",
        "result_parameter": "result",
        "search_type": "similarity_eval",
    }

    """Pptx_Extraction_Tool
    ---------------------
    Purpose:
        A comprehensive PowerPoint document extraction tool that extracts content from PowerPoint files (.pptx, .ppt) with support for multiple slides, embedded media extraction (images, audio, video), speaker notes, and LLM-friendly output formatting. Provides detailed metadata about slide structure, content types, processing engine used, and extracted media files.

    Core Capabilities:
        - Extracts text content from PowerPoint slides and titles
        - Extracts speaker notes from presentations
        - Extracts embedded images and media files
        - Analyzes presentation structure and slide layouts

    Intended Use:
        Use this tool when you need to extract content from PowerPoint presentations, including embedded media extraction, speaker notes extraction, and presentation structure analysis.

    Limitations:
        - May not handle complex PowerPoint presentations or embedded media
        - Large presentations may consume significant memory
        - Media extraction is limited to embedded content
        - Processing time scales with presentation size and complexity
        - Some advanced animations and effects may not be captured
        - Comments extraction may be limited
    """

    def __init__(self) -> None:
        super().__init__(
            type='function',
            name="Pptx_Extraction_Tool",
            description="""A comprehensive PowerPoint document extraction tool that extracts content from PowerPoint files (.pptx, .ppt) with support for multiple slides, embedded media extraction (images, audio, video), speaker notes extraction, and LLM-friendly output in multiple formats. Provides detailed metadata about slide structure, content types, processing engine used, and extracted media files. CAPABILITIES: Extracts text content from PowerPoint slides and titles, extracts speaker notes from presentations, extracts embedded images and media files, analyzes presentation structure and slide layouts. SYNONYMS: PowerPoint extractor, presentation parser, slide content extractor, PPTX reader, presentation analyzer, slide text extractor, PowerPoint content extractor, presentation content tool, slide parser, PowerPoint reader. EXAMPLES: 'Extract content from this PowerPoint presentation', 'Get speaker notes from these slides', 'Extract images from this presentation file'.""",
            parameters={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Path to the PowerPoint document file to extract content from"
                    },
                    "output_format": {
                        "type": "string",
                        "description": "Output format: 'markdown', 'json', 'html', or 'text' (default: json)",
                        "enum": ["markdown", "json", "html", "text"]
                    },
                    "extract_images": {
                        "type": "boolean",
                        "description": "Whether to extract and save embedded images from the presentation (default: True)"
                    },
                    "extract_notes": {
                        "type": "boolean",
                        "description": "Whether to extract speaker notes (default: True)"
                    },
                    "include_structure": {
                        "type": "boolean",
                        "description": "Whether to include presentation structure information (default: True)"
                    }
                },
                "required": ["file_path"],
                "additionalProperties": False,
            },
            strict=False,
            category="document_processing",
            tags=["powerpoint_extraction", "document_processing", "slide_extraction", "presentation_parser", "speaker_notes", "media_extraction", "slide_analyzer", "presentation_tool", "pptx_reader", "content_extraction"],
            limitation="Some complex PowerPoint features may not be fully preserved, large presentations may consume significant memory, media extraction limited to embedded content, processing time scales with presentation size and complexity, some advanced animations and effects may not be captured, comments extraction may be limited",
            agent_type="File_Extraction-Agent",
            accuracy= self.find_accuracy(os.path.join(os.path.dirname(__file__), 'test_result.json')),
            demo_commands= {
                "command": "reponse = tool.run(file_path='test.pptx', output_format='markdown', extract_images=True, extract_notes=True, include_structure=True)",
                "description": "Extract content from test.pptx in markdown format, extract images, extract notes, include structure"
            }
        )
        self.workspace = Path(os.getcwd())
        self._media_output_dir = self.workspace / "extracted_media"

        self.supported_extensions = {".pptx", ".ppt"}
    def embed_tool(self):
        return super().embed_tool()
    
    def get_metadata(self):
        """
        Returns the metadata for the Pptx_Extraction_Tool.

        Returns:
            dict: A dictionary containing the tool's metadata.
        """
        metadata = super().get_metadata()
        return metadata
    def _extract_images_from_pptx(self, file_path: Path, file_stem: str) -> list[dict[str, str]]:
        """Extract embedded images from PPTX file.

        Args:
            file_path: Path to the PPTX file
            file_stem: Base name for saving files

        Returns:
            List of dictionaries containing image file paths and metadata
        """
        saved_media = []
        self._media_output_dir.mkdir(exist_ok=True)

        try:
            # PPTX files are ZIP archives
            with zipfile.ZipFile(file_path, "r") as zip_file:
                # Find media files in the archive
                media_files = [f for f in zip_file.namelist() if f.startswith("ppt/media/")]

                for idx, media_file in enumerate(media_files):
                    try:
                        # Extract file extension
                        original_ext = Path(media_file).suffix
                        if not original_ext:
                            original_ext = ".png"  # Default extension

                        # Generate unique filename
                        media_filename = f"{file_stem}_media_{idx}{original_ext}"
                        media_path = self._media_output_dir / media_filename

                        # Extract and save media file
                        with zip_file.open(media_file) as source:
                            with open(media_path, "wb") as target:
                                target.write(source.read())

                        # Determine media type
                        media_type = (
                            "image" if original_ext.lower() in {".png", ".jpg", ".jpeg", ".gif", ".bmp"} else "media"
                        )

                        saved_media.append(
                            {
                                "type": media_type,
                                "path": str(media_path),
                                "filename": media_filename,
                                "original_path": media_file,
                            }
                        )


                    except Exception as e:
                        raise ValueError(f"Failed to extract media {media_file}: {e}")

        except Exception as e:
            raise ValueError(f"Failed to extract media from PPTX: {e}")

        return saved_media

    def _extract_slide_structure(self, presentation: PresentationType) -> dict[str, Any]:
        """Extract presentation structure information.

        Args:
            presentation: python-pptx Presentation object

        Returns:
            Dictionary containing structure metadata
        """
        structure = {
            "slide_count": len(presentation.slides),
            "slide_layouts": [],
            "has_notes": False,
            "has_comments": False,
            "slide_sizes": None,
        }

        # Get slide size information
        if hasattr(presentation.slide_width, "inches") and hasattr(presentation.slide_height, "inches"):
            structure["slide_sizes"] = {
                "width_inches": presentation.slide_width.inches,
                "height_inches": presentation.slide_height.inches,
            }

        # Analyze slide layouts and content
        for slide_idx, slide in enumerate(presentation.slides):
            layout_info = {
                "slide_index": slide_idx,
                "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, "name") else "Unknown",
                "shape_count": len(slide.shapes),
                "has_title": False,
                "has_content": False,
            }

            # Check for title and content
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    try:
                        if (
                            hasattr(shape, "placeholder_format")
                            and shape.placeholder_format is not None
                            and "title" in str(shape.placeholder_format.type).lower()
                        ):
                            layout_info["has_title"] = True
                        else:
                            layout_info["has_content"] = True
                    except (ValueError, AttributeError):
                        # Shape is not a placeholder or has no placeholder format
                        layout_info["has_content"] = True

            # Check for notes
            if hasattr(slide, "notes_slide") and slide.notes_slide:
                if hasattr(slide.notes_slide, "notes_text_frame") and slide.notes_slide.notes_text_frame.text.strip():
                    structure["has_notes"] = True

            structure["slide_layouts"].append(layout_info)

        return structure

    def _extract_content_from_pptx(self, file_path: Path, extract_notes: bool = True) -> dict[str, Any]:
        """Extract content from PPTX file using python-pptx.

        Args:
            file_path: Path to the PPTX file
            extract_notes: Whether to extract speaker notes

        Returns:
            Dictionary containing extracted content and metadata
        """
        start_time = time.time()

        try:
            # Load the presentation
            presentation = Presentation(str(file_path))

            # Extract slide content
            slides_content = []
            total_text_length = 0

            for slide_idx, slide in enumerate(presentation.slides):
                slide_data = {
                    "slide_number": slide_idx + 1,
                    "title": "",
                    "content": [],
                    "notes": "",
                    "shapes_count": len(slide.shapes),
                }

                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content = shape.text.strip()

                        # Try to identify if this is a title
                        try:
                            if (
                                hasattr(shape, "placeholder_format")
                                and shape.placeholder_format is not None
                                and "title" in str(shape.placeholder_format.type).lower()
                            ):
                                slide_data["title"] = text_content
                            else:
                                slide_data["content"].append(text_content)
                        except (ValueError, AttributeError):
                            # Shape is not a placeholder or has no placeholder format
                            slide_data["content"].append(text_content)

                        total_text_length += len(text_content)

                    # Extract text from tables if present
                    if hasattr(shape, "table"):
                        table_text = []
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                table_text.append(" | ".join(row_text))
                        if table_text:
                            slide_data["content"].append("\n".join(table_text))

                # Extract speaker notes if requested
                if extract_notes and hasattr(slide, "notes_slide"):
                    try:
                        if (
                            hasattr(slide.notes_slide, "notes_text_frame")
                            and slide.notes_slide.notes_text_frame.text.strip()
                        ):
                            slide_data["notes"] = slide.notes_slide.notes_text_frame.text.strip()
                            total_text_length += len(slide_data["notes"])
                    except Exception as e:
                        raise ValueError(f"Failed to extract notes from slide {slide_idx + 1}: {e}")

                slides_content.append(slide_data)

            # Extract presentation structure
            structure = self._extract_slide_structure(presentation)

            processing_time = time.time() - start_time

            return {
                "slides": slides_content,
                "structure": structure,
                "processing_time": processing_time,
                "total_text_length": total_text_length,
                "slide_count": len(slides_content),
            }

        except Exception as e:
            raise ValueError(f"Failed to extract content from PPTX: {e}")

    def _format_content_for_llm(
        self, extraction_result: dict[str, Any], output_format: str, include_structure: bool = True
    ) -> str:
        """Format extracted PPTX content to be LLM-friendly.

        Args:
            extraction_result: Dictionary containing extracted content
            output_format: Desired output format
            include_structure: Whether to include presentation structure information

        Returns:
            Formatted content string
        """
        slides = extraction_result["slides"]
        structure = extraction_result["structure"]

        if output_format.lower() == "markdown":
            content_parts = []

            if include_structure:
                content_parts.append("# Presentation Overview")
                content_parts.append(f"- **Total Slides**: {structure['slide_count']}")
                if structure.get("slide_sizes"):
                    sizes = structure["slide_sizes"]
                    content_parts.append(
                        f'- **Slide Size**: {sizes["width_inches"]:.1f}" × {sizes["height_inches"]:.1f}"'
                    )
                content_parts.append(f"- **Has Speaker Notes**: {'Yes' if structure['has_notes'] else 'No'}")
                content_parts.append("")

            # Format each slide
            for slide in slides:
                content_parts.append(f"## Slide {slide['slide_number']}")

                if slide["title"]:
                    content_parts.append(f"**Title**: {slide['title']}")
                    content_parts.append("")

                if slide["content"]:
                    content_parts.append("**Content**:")
                    for content_item in slide["content"]:
                        # Format multi-line content properly
                        for line in content_item.split("\n"):
                            if line.strip():
                                content_parts.append(f"- {line.strip()}")
                    content_parts.append("")

                if slide["notes"]:
                    content_parts.append("**Speaker Notes**:")
                    content_parts.append(slide["notes"])
                    content_parts.append("")

                content_parts.append("---")
                content_parts.append("")

            return "\n".join(content_parts)

        elif output_format.lower() == "json":
            return json.dumps(
                {
                    "presentation_structure": structure if include_structure else None,
                    "slides": slides,
                    "metadata": {
                        "total_slides": len(slides),
                        "total_text_length": extraction_result["total_text_length"],
                        "processing_time": extraction_result["processing_time"],
                    },
                },
                indent=2,
            )

        elif output_format.lower() == "html":
            html_parts = ["<div class='presentation'>"]

            if include_structure:
                html_parts.append("<div class='presentation-overview'>")
                html_parts.append("<h1>Presentation Overview</h1>")
                html_parts.append(f"<p><strong>Total Slides:</strong> {structure['slide_count']}</p>")
                if structure.get("slide_sizes"):
                    sizes = structure["slide_sizes"]
                    html_parts.append(
                        "<p>"
                        "<strong>Slide Size:</strong> "
                        f"{sizes['width_inches']:.1f} × {sizes['height_inches']:.1f}"
                        "</p>"
                    )
                html_parts.append(
                    f"<p><strong>Has Speaker Notes:</strong> {'Yes' if structure['has_notes'] else 'No'}</p>"
                )
                html_parts.append("</div>")

            for slide in slides:
                html_parts.append(f"<div class='slide' data-slide='{slide['slide_number']}'>")
                html_parts.append(f"<h2>Slide {slide['slide_number']}</h2>")

                if slide["title"]:
                    html_parts.append(f"<h3>{slide['title']}</h3>")

                if slide["content"]:
                    html_parts.append("<div class='slide-content'>")
                    for content_item in slide["content"]:
                        html_parts.append(f"<p>{content_item.replace(chr(10), '<br>')}</p>")
                    html_parts.append("</div>")

                if slide["notes"]:
                    html_parts.append(
                        "<div class='speaker-notes'>"
                        "<strong>Speaker Notes:</strong>"
                        f"<br>{slide['notes'].replace(chr(10), '<br>')}"
                        "</div>"
                    )

                html_parts.append("</div>")

            html_parts.append("</div>")
            return "\n".join(html_parts)

        else:  # Plain text
            text_parts = []

            if include_structure:
                text_parts.append("PRESENTATION OVERVIEW")
                text_parts.append(f"Total Slides: {structure['slide_count']}")
                if structure.get("slide_sizes"):
                    sizes = structure["slide_sizes"]
                    text_parts.append(f'Slide Size: {sizes["width_inches"]:.1f}" × {sizes["height_inches"]:.1f}"')
                text_parts.append(f"Has Speaker Notes: {'Yes' if structure['has_notes'] else 'No'}")
                text_parts.append("\n" + "=" * 50 + "\n")

            for slide in slides:
                text_parts.append(f"SLIDE {slide['slide_number']}")

                if slide["title"]:
                    text_parts.append(f"Title: {slide['title']}")

                if slide["content"]:
                    text_parts.append("Content:")
                    for content_item in slide["content"]:
                        text_parts.append(f"  {content_item}")

                if slide["notes"]:
                    text_parts.append(f"Speaker Notes: {slide['notes']}")

                text_parts.append("\n" + "-" * 30 + "\n")

            return "\n".join(text_parts)

    def run(
        self,
        file_path: str = Field(description="Path to the PPTX/PPT presentation file to extract content from"),
        output_format: Literal["markdown", "json", "html", "text"] = Field(
            default="json", description="Output format: 'markdown', 'json', 'html', or 'text'"
        ),
        extract_images: bool = Field(
            default=True, description="Whether to extract and save images from the presentation"
        ),
        extract_notes: bool = Field(default=True, description="Whether to extract speaker notes"),
        include_structure: bool = Field(
            default=True, description="Whether to include presentation structure information"
        ),
    ) :
        """Extract content from PPTX/PPT presentations using python-pptx.

        This tool provides comprehensive PowerPoint presentation content extraction with support for:
        - PPTX and PPT files
        - Text extraction from slides, titles, and content
        - Speaker notes extraction
        - Image and media extraction
        - Presentation structure analysis
        - Multiple output formats (Markdown, JSON, HTML, Text)
        - LLM-optimized formatting

        Args:
            file_path: Path to the PPTX/PPT presentation file
            output_format: Desired output format
            extract_images: Extract embedded media files
            extract_notes: Extract speaker notes
            include_structure: Include presentation structure info

        Returns:
            ActionResponse with extracted content, metadata, and media file paths
        """
        try:
            # Handle FieldInfo objects from pydantic
            if isinstance(file_path, FieldInfo):
                file_path = file_path.default
            if isinstance(output_format, FieldInfo):
                output_format = output_format.default
            if isinstance(extract_images, FieldInfo):
                extract_images = extract_images.default
            if isinstance(extract_notes, FieldInfo):
                extract_notes = extract_notes.default
            if isinstance(include_structure, FieldInfo):
                include_structure = include_structure.default

            # Validate input file
            file_path: Path = self._validate_file_path(file_path)
            print(f"Processing PPTX presentation: {file_path.name}")

            # Extract embedded media if requested
            saved_media = []
            if extract_images and file_path.suffix.lower() == ".pptx":
                saved_media = self._extract_images_from_pptx(file_path, file_path.stem)

            # Extract presentation content
            extraction_result = self._extract_content_from_pptx(file_path, extract_notes=extract_notes)

            # Format content for LLM consumption
            formatted_content = self._format_content_for_llm(
                extraction_result, output_format, include_structure=include_structure
            )

            # Prepare metadata
            file_stats = file_path.stat()
            document_metadata = DocumentMetadata(
                file_name=file_path.name,
                file_size=file_stats.st_size,
                file_type=file_path.suffix.lower(),
                absolute_path=str(file_path.absolute()),
                page_count=extraction_result["slide_count"],  # Use slide count as page count
                processing_time=extraction_result["processing_time"],
                extracted_images=[media["path"] for media in saved_media if media["type"] == "image"],
                extracted_media=saved_media,
                output_format=output_format,
                llm_enhanced=False,
                ocr_applied=False,
            )

            # Add PPTX-specific metadata
            pptx_metadata = {
                "slide_count": extraction_result["slide_count"],
                "total_text_length": extraction_result["total_text_length"],
                "has_speaker_notes": extraction_result["structure"]["has_notes"],
                "slide_layouts": extraction_result["structure"]["slide_layouts"],
                "presentation_size": extraction_result["structure"].get("slide_sizes"),
                "media_files_count": len(saved_media),
            }

            # Merge metadata
            final_metadata = {**document_metadata.model_dump(), **pptx_metadata}

            return {'result': formatted_content, 'success': True, 'metadata': final_metadata}

        except FileNotFoundError as e:
            return {"error": f"File not found: {str(e)}", "success": False, "error_type": "file_not_found", "traceback": traceback.format_exc()}
        except ValueError as e:
            return {"error": f"Invalid input: {str(e)}", "success": False, "error_type": "invalid_input", "traceback": traceback.format_exc()}
        except Exception as e:
            return {"error": f"PPTX extraction failed: {str(e)}", "success": False, "error_type": "pptx_extraction_failed", "traceback": traceback.format_exc()}
    def _validate_file_path(self, file_path: str) -> Path:
        """Validate and resolve file path. Rely on the predefined supported_extensions class variable.

        Args:
            file_path: Path to the document or media file

        Returns:
            Resolved Path object

        Raises:
            FileNotFoundError: If file doesn't exist
            ValueError: If file type is not supported
        """
        path = Path(file_path).expanduser()
        if not path.is_absolute():
            path = self.workspace / path

        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        if path.suffix.lower() not in self.supported_extensions:
            raise ValueError(
                f"Unsupported file type: {path.suffix}. Supported types: {', '.join(self.supported_extensions)}"
            )

        return path

    def test(self, tool_test: str="pptx_extraction", file_location: str="pptx_extraction", result_parameter: str="result", search_type: str="exact_match"):
        return super().test(tool_test=tool_test, file_location=file_location, result_parameter=result_parameter, search_type=search_type)
# Example usage and entry point
if __name__ == "__main__":

    # Initialize and run the PPTX extraction service
    try:
        tool = Pptx_Extraction_Tool()
        tool.embed_tool()
        tool.test(tool_test="pptx_extraction", file_location="pptx_extraction", result_parameter="result", search_type="exact_match")
    except Exception as e:
        print(f"An error occurred: {e}: {traceback.format_exc()}")
